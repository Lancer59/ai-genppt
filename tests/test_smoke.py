"""
Smoke test — validates schema parsing, theme loading, and builder assembly
WITHOUT needing an LLM key or network access.
"""

import json
from pathlib import Path

import pytest
from pptx import Presentation

from genppt.llm.schemas import SlideOutline, SlideSpec, StatItem
from genppt.renderer.builder import PresentationBuilder
from genppt.renderer.themes import THEMES, get_theme


# ── Sample outline fixture ───────────────────────────────────────────────────

@pytest.fixture
def sample_outline() -> SlideOutline:
    return SlideOutline(
        presentation_title="Test Presentation",
        theme="dark",
        total_slides=4,
        slides=[
            SlideSpec(
                layout="title",
                title="Welcome to the Future",
                subtitle="An AI-generated deck",
                speaker_notes="Introduce yourself and the topic.",
                icon_names=["rocket"],
            ),
            SlideSpec(
                layout="content",
                title="Key Insights",
                bullets=["Insight one", "Insight two", "Insight three"],
                image_query="technology abstract",
                speaker_notes="Walk through each key insight.",
                icon_names=["bar-chart", "lightning-charge"],
            ),
            SlideSpec(
                layout="stats_callout",
                title="By the Numbers",
                stats=[
                    StatItem(value="98%", label="Accuracy"),
                    StatItem(value="10x", label="Faster"),
                    StatItem(value="$4.2B", label="Market Size"),
                ],
                speaker_notes="Highlight the impact numbers.",
            ),
            SlideSpec(
                layout="closing",
                title="Thank You",
                subtitle="Questions?",
                speaker_notes="Open for Q&A.",
            ),
        ],
    )


# ── Schema tests ─────────────────────────────────────────────────────────────

class TestSchemas:
    def test_slide_outline_roundtrip(self, sample_outline):
        """Outline can serialize to JSON and parse back."""
        data = sample_outline.model_dump()
        parsed = SlideOutline(**data)
        assert parsed.presentation_title == "Test Presentation"
        assert len(parsed.slides) == 4

    def test_slide_spec_layout_validation(self):
        """Invalid layout should raise validation error."""
        with pytest.raises(Exception):
            SlideSpec(layout="nonexistent", title="Bad")

    def test_stat_item(self):
        s = StatItem(value="42%", label="Hit Rate")
        assert s.value == "42%"


# ── Theme tests ──────────────────────────────────────────────────────────────

class TestThemes:
    def test_all_themes_load(self):
        assert len(THEMES) >= 3
        for name in ("dark", "corporate", "minimalist"):
            theme = get_theme(name)
            assert theme.name == name
            assert len(theme.bg_color) == 3

    def test_invalid_theme_raises(self):
        with pytest.raises(ValueError):
            get_theme("nonexistent")


# ── Builder tests ────────────────────────────────────────────────────────────

class TestBuilder:
    def test_build_without_assets(self, sample_outline, tmp_path):
        """Builder produces a valid .pptx even with zero images/icons."""
        output = str(tmp_path / "test_deck.pptx")
        builder = PresentationBuilder(theme_name="dark")
        result = builder.build(
            outline=sample_outline,
            image_paths={},
            icon_paths={},
            output_path=output,
        )
        assert Path(result).exists()
        # Verify the file is a valid PPTX
        prs = Presentation(result)
        assert len(prs.slides) == 4

    def test_build_all_themes(self, sample_outline, tmp_path):
        """Each theme produces a valid .pptx."""
        for theme_name in ("dark", "corporate", "minimalist"):
            output = str(tmp_path / f"{theme_name}.pptx")
            builder = PresentationBuilder(theme_name=theme_name)
            result = builder.build(
                outline=sample_outline,
                image_paths={},
                icon_paths={},
                output_path=output,
            )
            assert Path(result).exists()

    def test_speaker_notes_present(self, sample_outline, tmp_path):
        """Speaker notes from the outline end up in the PPTX."""
        output = str(tmp_path / "notes_test.pptx")
        builder = PresentationBuilder(theme_name="corporate")
        builder.build(
            outline=sample_outline,
            image_paths={},
            icon_paths={},
            output_path=output,
        )
        prs = Presentation(output)
        slide_1 = prs.slides[0]
        notes = slide_1.notes_slide.notes_text_frame.text
        assert "Introduce yourself" in notes
