"""
PresentationBuilder — orchestrates python-pptx to assemble the final .pptx.

Responsibilities:
  - Load a blank/template presentation
  - For each slide in the outline, add a blank slide and call the correct layout renderer
  - Inject speaker notes
  - Save the output file
"""

from pathlib import Path
from typing import Dict, Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from genppt.llm.schemas import SlideOutline
from genppt.renderer.layouts import LAYOUT_RENDERERS
from genppt.renderer.themes import get_theme
from genppt.renderer.themes.base import Theme
from genppt.utils.logger import get_logger

logger = get_logger(__name__)

# Widescreen 16:9
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


class PresentationBuilder:
    def __init__(
        self,
        theme_name: str = "dark",
        template_path: Optional[str] = None,
    ) -> None:
        self.theme: Theme = get_theme(theme_name)
        self.template_path = template_path

    def _make_prs(self) -> Presentation:
        if self.template_path and Path(self.template_path).exists():
            prs = Presentation(self.template_path)
        else:
            prs = Presentation()
            prs.slide_width = SLIDE_WIDTH
            prs.slide_height = SLIDE_HEIGHT
        return prs

    def _blank_slide(self, prs: Presentation):
        """Add a completely blank slide (no placeholders)."""
        blank_layout = prs.slide_layouts[6]  # index 6 = Blank in standard template
        return prs.slides.add_slide(blank_layout)

    def _add_speaker_notes(self, slide, notes_text: str) -> None:
        if not notes_text:
            return
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    def build(
        self,
        outline: SlideOutline,
        image_paths: Dict[int, str],
        icon_paths: Dict[str, str],
        output_path: str,
    ) -> str:
        prs = self._make_prs()
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)

        logger.info(f"Rendering {len(outline.slides)} slides (theme={self.theme.name})")

        for i, spec in enumerate(outline.slides):
            slide = self._blank_slide(prs)

            # Resolve renderer
            renderer_fn = LAYOUT_RENDERERS.get(spec.layout, LAYOUT_RENDERERS["content"])

            # Image for this slide
            img_path = image_paths.get(i)

            try:
                renderer_fn(
                    slide=slide,
                    spec=spec,
                    theme=self.theme,
                    image_path=img_path,
                    icon_paths=icon_paths,
                )
            except Exception as e:
                logger.warning(f"Slide {i + 1} render error ({spec.layout}): {e}")

            # Speaker notes
            self._add_speaker_notes(slide, spec.speaker_notes)

        prs.save(str(out))
        logger.info(f"Saved -> {out}")
        return str(out)
