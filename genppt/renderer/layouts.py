"""
Slide layout implementations.

Each layout function receives:
  - slide         : pptx Slide object (blank layout, background already set)
  - spec          : SlideSpec — the LLM-generated content
  - theme         : Theme dataclass
  - image_path    : Optional local file path for a photo
  - icon_paths    : Dict[icon_name, local_png_path]

All measurements use python-pptx Inches/Pt helpers.
"""

import io
from pathlib import Path
from typing import Dict, List, Optional

from PIL import Image as PILImage
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from genppt.llm.schemas import SlideSpec
from genppt.renderer.themes.base import RGB, Theme

# Slide dimensions (widescreen 16:9)
SW = Inches(13.33)   # slide width
SH = Inches(7.5)     # slide height


# ── Helpers ──────────────────────────────────────────────────────────────────

def _rgb(color: RGB) -> RGBColor:
    return RGBColor(*color)


def _set_bg(slide, color: RGB) -> None:
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_textbox(slide, left, top, width, height, text: str,
                  font_name: str, font_size: int, color: RGB,
                  bold: bool = False, align=PP_ALIGN.LEFT,
                  wrap: bool = True) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = _rgb(color)
    run.font.bold = bold


def _add_bullet_textbox(slide, left, top, width, height,
                         bullets: List[str], font_name: str,
                         font_size: int, color: RGB) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        for run in p.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = _rgb(color)


def _add_image(slide, img_path: str, left, top, width, height) -> None:
    try:
        slide.shapes.add_picture(img_path, left, top, width, height)
    except Exception:
        pass


def _add_accent_bar(slide, theme: Theme, height_frac: float = 0.008) -> None:
    """Thin colored bar at top of slide."""
    bar_h = SH * height_frac
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SW, bar_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(theme.accent_color)
    shape.line.fill.background()


def _add_icons_row(slide, icon_paths: Dict[str, str],
                   icon_names: Optional[List[str]],
                   left, top, size=Inches(0.45)) -> None:
    if not icon_names or not icon_paths:
        return
    x = left
    for name in icon_names[:3]:
        path = icon_paths.get(name)
        if path and Path(path).exists():
            try:
                slide.shapes.add_picture(path, x, top, size, size)
            except Exception:
                pass
            x += size + Inches(0.1)


# ── Layout functions ─────────────────────────────────────────────────────────

def render_title(slide, spec: SlideSpec, theme: Theme,
                 image_path: Optional[str], icon_paths: Dict[str, str]) -> None:
    _set_bg(slide, theme.slide_bg)
    _add_accent_bar(slide, theme, 0.012)

    # Background image with dark overlay if available
    if image_path:
        _add_image(slide, image_path, 0, 0, SW, SH)
        # Semi-transparent overlay via a rectangle
        overlay = slide.shapes.add_shape(1, 0, 0, SW, SH)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = _rgb(theme.slide_bg)
        from pptx.util import Pt as _Pt
        overlay.fill.fore_color.theme_color  # touch to keep solid
        # Set transparency via XML
        from lxml import etree
        sp_tree = overlay._element
        solidFill = sp_tree.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if solidFill is not None:
            srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgb is not None:
                alpha = etree.SubElement(
                    srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
                )
                alpha.set('val', '75000')  # ~75% opacity

    # Title
    _add_textbox(
        slide, Inches(1.2), Inches(2.2), Inches(10.9), Inches(2.2),
        spec.title, theme.heading_font, theme.title_size,
        theme.title_color, bold=True, align=PP_ALIGN.CENTER
    )
    # Subtitle
    if spec.subtitle:
        _add_textbox(
            slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(1.2),
            spec.subtitle, theme.body_font, theme.subtitle_size,
            theme.subtitle_color, align=PP_ALIGN.CENTER
        )
    # Icons bottom-right
    _add_icons_row(slide, icon_paths, spec.icon_names,
                   Inches(11.0), Inches(6.6), Inches(0.4))


def render_section_header(slide, spec: SlideSpec, theme: Theme,
                           image_path: Optional[str], icon_paths: Dict[str, str]) -> None:
    _set_bg(slide, theme.accent_color)
    _add_textbox(
        slide, Inches(1.0), Inches(2.8), Inches(11.3), Inches(1.8),
        spec.title, theme.heading_font, theme.title_size,
        (255, 255, 255), bold=True, align=PP_ALIGN.CENTER
    )
    if spec.subtitle:
        _add_textbox(
            slide, Inches(2.0), Inches(4.7), Inches(9.3), Inches(1.0),
            spec.subtitle, theme.body_font, theme.subtitle_size,
            (220, 220, 255), align=PP_ALIGN.CENTER
        )
    _add_icons_row(slide, icon_paths, spec.icon_names,
                   Inches(6.0), Inches(1.0), Inches(0.5))


def render_content(slide, spec: SlideSpec, theme: Theme,
                   image_path: Optional[str], icon_paths: Dict[str, str]) -> None:
    _set_bg(slide, theme.slide_bg)
    _add_accent_bar(slide, theme, 0.008)

    has_image = bool(image_path)

    # Title
    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0),
        spec.title, theme.heading_font, theme.title_size - 6,
        theme.title_color, bold=True
    )

    # Accent line under title
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(theme.accent_color)
    line.line.fill.background()

    # Icons row below accent line
    _add_icons_row(slide, icon_paths, spec.icon_names,
                   Inches(0.5), Inches(1.5), Inches(0.4))

    bullet_left = Inches(0.5)
    bullet_width = Inches(7.5) if has_image else Inches(12.3)

    if spec.bullets:
        _add_bullet_textbox(
            slide, bullet_left, Inches(2.0), bullet_width, Inches(5.0),
            spec.bullets, theme.body_font, theme.body_size, theme.body_color
        )

    if has_image:
        _add_image(slide, image_path, Inches(8.2), Inches(1.4), Inches(4.7), Inches(5.7))


def render_two_column(slide, spec: SlideSpec, theme: Theme,
                       image_path: Optional[str], icon_paths: Dict[str, str]) -> None:
    _set_bg(slide, theme.slide_bg)
    _add_accent_bar(slide, theme, 0.008)

    # Main title
    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0),
        spec.title, theme.heading_font, theme.title_size - 6,
        theme.title_color, bold=True
    )

    # Divider line
    div = slide.shapes.add_shape(1, Inches(6.5), Inches(1.3), Inches(0.04), Inches(5.8))
    div.fill.solid()
    div.fill.fore_color.rgb = _rgb(theme.accent_color)
    div.line.fill.background()

    # Left column
    if spec.left_title:
        _add_textbox(
            slide, Inches(0.4), Inches(1.4), Inches(5.8), Inches(0.7),
            spec.left_title, theme.heading_font, theme.body_size + 2,
            theme.accent_color, bold=True
        )
    if spec.left_bullets:
        _add_bullet_textbox(
            slide, Inches(0.4), Inches(2.2), Inches(5.8), Inches(4.8),
            spec.left_bullets, theme.body_font, theme.body_size, theme.body_color
        )

    # Right column
    if spec.right_title:
        _add_textbox(
            slide, Inches(6.7), Inches(1.4), Inches(6.2), Inches(0.7),
            spec.right_title, theme.heading_font, theme.body_size + 2,
            theme.accent_color, bold=True
        )
    if spec.right_bullets:
        _add_bullet_textbox(
            slide, Inches(6.7), Inches(2.2), Inches(6.2), Inches(4.8),
            spec.right_bullets, theme.body_font, theme.body_size, theme.body_color
        )


def render_full_bleed_image(slide, spec: SlideSpec, theme: Theme,
                             image_path: Optional[str], icon_paths: Dict[str, str]) -> None:
    _set_bg(slide, theme.slide_bg)
    if image_path:
        _add_image(slide, image_path, 0, 0, SW, SH)

    # Dark gradient overlay at bottom for legibility
    overlay = slide.shapes.add_shape(1, 0, Inches(3.5), SW, Inches(4.0))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    from lxml import etree
    sp = overlay._element
    solidFill = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solidFill is not None:
        srgb = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha = etree.SubElement(
                srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
            )
            alpha.set('val', '60000')
    overlay.line.fill.background()

    _add_textbox(
        slide, Inches(0.7), Inches(4.0), Inches(11.9), Inches(1.6),
        spec.title, theme.heading_font, theme.title_size,
        (255, 255, 255), bold=True
    )
    if spec.subtitle:
        _add_textbox(
            slide, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.0),
            spec.subtitle, theme.body_font, theme.subtitle_size,
            (210, 210, 210)
        )


def render_stats_callout(slide, spec: SlideSpec, theme: Theme,
                          image_path: Optional[str], icon_paths: Dict[str, str]) -> None:
    _set_bg(slide, theme.slide_bg)
    _add_accent_bar(slide, theme, 0.012)

    _add_textbox(
        slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(1.0),
        spec.title, theme.heading_font, theme.title_size - 6,
        theme.title_color, bold=True, align=PP_ALIGN.CENTER
    )

    stats = spec.stats or []
    n = len(stats)
    if n == 0:
        return

    col_w = Inches(13.33 / n)
    for i, stat in enumerate(stats):
        left = col_w * i
        # Stat value
        _add_textbox(
            slide, left, Inches(2.2), col_w, Inches(2.0),
            stat.value, theme.heading_font, 54,
            theme.stat_color, bold=True, align=PP_ALIGN.CENTER
        )
        # Stat label
        _add_textbox(
            slide, left, Inches(4.3), col_w, Inches(1.0),
            stat.label, theme.body_font, theme.body_size,
            theme.body_color, align=PP_ALIGN.CENTER
        )
        # Thin accent underline
        ul = slide.shapes.add_shape(
            1,
            left + col_w * 0.2, Inches(4.2),
            col_w * 0.6, Inches(0.04)
        )
        ul.fill.solid()
        ul.fill.fore_color.rgb = _rgb(theme.accent_color)
        ul.line.fill.background()

    _add_icons_row(slide, icon_paths, spec.icon_names,
                   Inches(0.5), Inches(5.7), Inches(0.5))


def render_closing(slide, spec: SlideSpec, theme: Theme,
                   image_path: Optional[str], icon_paths: Dict[str, str]) -> None:
    _set_bg(slide, theme.accent_color)

    _add_textbox(
        slide, Inches(1.0), Inches(2.3), Inches(11.3), Inches(2.0),
        spec.title, theme.heading_font, theme.title_size + 4,
        (255, 255, 255), bold=True, align=PP_ALIGN.CENTER
    )
    if spec.subtitle:
        _add_textbox(
            slide, Inches(2.0), Inches(4.5), Inches(9.3), Inches(1.2),
            spec.subtitle, theme.body_font, theme.subtitle_size,
            (220, 220, 255), align=PP_ALIGN.CENTER
        )
    _add_icons_row(slide, icon_paths, spec.icon_names,
                   Inches(6.1), Inches(6.2), Inches(0.45))


# ── Dispatch table ───────────────────────────────────────────────────────────
LAYOUT_RENDERERS = {
    "title": render_title,
    "section_header": render_section_header,
    "content": render_content,
    "two_column": render_two_column,
    "full_bleed_image": render_full_bleed_image,
    "stats_callout": render_stats_callout,
    "closing": render_closing,
}
